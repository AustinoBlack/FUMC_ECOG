from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import sys
import os
import datetime
import subprocess

# Function to convert a color name string to an RGBColor
def get_rgb_from_name(color_name):
    color_map = {
        "green": RGBColor(0, 255, 0),
        "white": RGBColor(255, 255, 255),
        "lightgray": RGBColor(211, 211, 211),
        "black": RGBColor(0, 0, 0),
        "beige": RGBColor(245, 245, 220),
        "blue": RGBColor(0, 112, 192),
        "red": RGBColor(255, 0, 0),
        "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0),
        "purple": RGBColor(128, 0, 128),
        "teal": RGBColor(0, 128, 128),
        "navy": RGBColor(0, 0, 128),
        "gray": RGBColor(128, 128, 128),
        "pink": RGBColor(255, 192, 203)
    }
    return color_map.get(color_name.lower(), RGBColor(0, 255, 0))  # default: green

def get_upcoming_sunday():
    today = datetime.date.today()
    days_until_sunday = (6 - today.weekday()) % 7
    upcoming_sunday = today + datetime.timedelta(days=days_until_sunday)
    return upcoming_sunday.strftime("%m-%d-%Y")

# Function to create the output folder
def create_output_folder():
    upcoming_sunday = get_upcoming_sunday()
    output_folder = os.path.join("outputs", upcoming_sunday)
    os.makedirs(output_folder, exist_ok=True)
    return output_folder

# Function to create background fill
def create_background(slide, slide_width, slide_height, background_color):
    shape = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = background_color

# Function to create the bottom bar
def create_bottom_bar(slide, slide_width, slide_height):
    textbox_img = 'assets/text_plate.png'
    bar_height = Inches(2.5)
    bar_top = slide_height - bar_height
    shape = slide.shapes.add_picture(textbox_img, Inches(1.5), bar_top, Inches(14.5), bar_height)
    return bar_top

# Function to insert icon on a slide
def insert_icon(slide, bar_top, icon_path):
    if icon_path and os.path.exists(icon_path):
        icon_left = Inches(-1.25)
        icon_top = bar_top
        slide.shapes.add_picture(icon_path, icon_left, icon_top, height=Inches(2.5))

# Function to insert text on a slide
def insert_text(slide, bar_top, slide_width, text, font_name):
    text_left = Inches(2)
    text_top = bar_top + Inches(0.25)
    textbox = slide.shapes.add_textbox(left=text_left, top=text_top, width=Inches(14.5), height=Inches(2.5))
    text_frame = textbox.text_frame
    text_frame.auto_size = None
    text_frame.word_wrap = True

    # Center Horizontally
    text = text_frame.paragraphs[0]
    text.alignment = PP_ALIGN.CENTER

    # Center Vertically
    text_fram.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(32)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(0, 0, 0)

# Function to extract text from a slide
def extract_text(slide):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
    return text.strip()

# Function to add background image
def create_background_image(slide, slide_width, slide_height, bg_image_path):
    slide.shapes.add_picture(bg_image_path, 0, 0, width=slide_width, height=slide_height)

# Function to process each slide of the presentation
def process_slide(new_ppt, new_slide, og_slide, font_name, icon_path, background_color, is_image_bg, bg_image_path):
    width = new_ppt.slide_width
    height = new_ppt.slide_height
    if is_image_bg == 1 and bg_image_path:
        create_background_image(new_slide, width, height, bg_image_path)
    else:
        create_background(new_slide, width, height, background_color)

    bar_top = create_bottom_bar(new_slide, width, height)
    ex_text = extract_text(og_slide)
    insert_icon(new_slide, bar_top, icon_path)
    insert_text(new_slide, bar_top, width, ex_text, font_name)

# Function to process the entire powerpoint
def process_pptx(input_file, output_folder, icon_path, background_color_name, font_name, is_image_bg, bg_image_path):
    output_file = os.path.join(output_folder, "output.pptx")
    print(f"Processing {input_file}...")
    print(f"Output will be saved as: {output_file}")

    prs = Presentation(input_file)
    new_ppt = Presentation()

    # Set slide width and height
    new_ppt.slide_width = Inches(16)
    new_ppt.slide_height = Inches(9)

    # Convert background color name to RGBColor
    background_color = get_rgb_from_name(background_color_name)

    for slide in prs.slides:
        layout = new_ppt.slide_layouts[6]  # 6 = blank slide
        new_slide = new_ppt.slides.add_slide(layout)
        process_slide(new_ppt, new_slide, slide, font_name, icon_path, background_color, is_image_bg, bg_image_path)

    new_ppt.save(output_file)
    print(f"PowerPoint saved as {output_file}")

# Function to generate preview images
def generate_preview(input_path, preview_path):
    os.makedirs(preview_path, exist_ok=True)

    try:
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'png', input_path, '--outdir', preview_path], check=True)
        print(f"Conversion complete, PNG files saved to {preview_path}")
    except subprocess.CalledProcessError as e:
        print(f"Error during conversion: {e}")
        return

    images = sorted([f for f in os.listdir(preview_path) if f.lower().endswith('.png')])
    print("Preview images ready.")
    return images

if __name__ == "__main__":
    input_pptx = "input.pptx"
    output_folder = create_output_folder()

    icon_path = "church.png"
    background_color_name = "green"
    font_name = "Arial"
    is_image_bg = 1
    bg_image_path = "church.png"

    process_pptx(input_pptx, output_folder, icon_path, background_color_name, font_name, is_image_bg, bg_image_path)

