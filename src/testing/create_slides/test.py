from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(16)  # Width in inches
prs.slide_height = Inches(9)  # Height in inches
blank_slide_layout = prs.slide_layouts[6]

# extract slide text
# get number of slides


#add "Title Slide"
def add_title_slide( content ):
    ''' adds a single slide of the title format to the presentation '''
    slide = prs.slides.add_slide(blank_slide_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 255, 0)

    logo_img = '../../assets/CrossFlame_WhiteRed.png'
    textbox_img = '../../assets/text_plate.png'
    
    # add text plate
    X = Inches(1.5)
    Y = Inches(6.5)
    height = Inches(2.5)
    width = Inches(14.5)
    textbox = slide.shapes.add_picture(textbox_img, X, Y, width=width, height=height)

    # add logo
    X = Inches(-1.25)
    Y = Inches(6.5)
    height = Inches(2.5)
    logo = slide.shapes.add_picture(logo_img, X, Y, height=height)

    # add text
    X = Inches(1.5)
    Y = Inches(6.5)
    height = Inches(2.5)
    width = Inches(14.5)
    text_box = slide.shapes.add_textbox(X, Y, width, height)
    text_frame = text_box.text_frame
    text_frame.auto_size = None
    text_frame.word_wrap = True
    text_frame.text = content

    # Center the text horizontally in the text box
    text = text_frame.paragraphs[0]  # Access the first paragraph (text content)
    text.alignment = PP_ALIGN.CENTER  # Center the text horizontally

    # Center the text vertically in the text box
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center the text vertically

    font = text.font
    font.name = 'Arial'  # Set font name
    font.size = Pt(40)   # Set font size to 24 pt
    font.bold = True     # Make the text bold (optional)


# add "Speech Slides"
def add_speech_slides( n, content ):
    ''' adds n slides of the speech format to the the presentation. Notes that the first of the speech slides will always be in the title format ''' 
    total = n
    i = 1
    
    while i <= total:
        slide = prs.slides.add_slide(blank_slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 255, 0)

        logo_img = '../../assets/CrossFlame_WhiteRed.png'
        textbox_img = '../../assets/text_plate.png'
    
        # add text plate
        X = Inches(1.5)
        Y = Inches(6.5)
        height = Inches(2.5)
        width = Inches(14.5)
        textbox = slide.shapes.add_picture(textbox_img, X, Y, width=width, height=height)

        # add logo
        X = Inches(-1.25)
        Y = Inches(6.5)
        height = Inches(2.5)
        logo = slide.shapes.add_picture(logo_img, X, Y, height=height)

        # add text
        X = Inches(1.5)
        Y = Inches(6.5)
        height = Inches(2.5)
        width = Inches(14.5)
        text_box = slide.shapes.add_textbox(X, Y, width, height)
        text_frame = text_box.text_frame
        text_frame.auto_size = None
        text_frame.word_wrap = True
        text_frame.text = content

        # Center the text horizontally in the text box
        text = text_frame.paragraphs[0]  # Access the first paragraph (text content)
        text.alignment = PP_ALIGN.CENTER  # Center the text horizontally

        # Center the text vertically in the text box
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center the text vertically

        font = text.font
        font.name = 'Arial'  # Set font name
        font.size = Pt(24)   # Set font size to 24 pt
        font.bold = True     # Make the text bold (optional)

        i += 1

title_content = "Title Slide: Placeholder Name"
group_content = "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat. Duis aute irure dolor in reprehenderit in voluptate."

add_title_slide(title_content)
add_speech_slides(3, group_content)

prs.save('test.pptx')
