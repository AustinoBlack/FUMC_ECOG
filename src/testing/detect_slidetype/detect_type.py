from pptx import Presentation

prs = Presentation("../../assets/test_slides.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

grouped_triggers = ['UMH', 'FWS', 'Scripture Reading', 'Prayer for Illumination', "The Lord’s Prayer"]
standalone_triggers = ['Passing of the Peace', 'Rev.', 'Prelude', 'Postlude', 'The Children’s Moment', 'Offering Our Gifts', 'Offertory', 'Sending Forth', 'Proclamation of God’s Word']

i=1
Flag = "Standalone"
for slide in prs.slides:
    data = []
    data.append(i)
    i+=1
    cleaned_text = ""
    for shape in slide.shapes:
        text = ""
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = ''.join( run.text.strip() ) + ' '
                if any( substr in text for substr in grouped_triggers) == True:
                    Flag = "Grouped"
                if any( substr in text for substr in standalone_triggers) == True:
                    Flag = "Standalone"
                if text != '':
                    cleaned_text += text
    data.append(Flag)
    data.append(cleaned_text)
    text_runs.append(data)
                     
for run in text_runs:
    print( run )

print( "Total Slides: " + str(i-1) )
