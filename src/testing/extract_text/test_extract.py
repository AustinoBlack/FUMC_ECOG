from pptx import Presentation

prs = Presentation("../../assets/test_slides.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

i=1
for slide in prs.slides:
    text_runs.append(i)
    i+=1
    cleaned_text = ""
    for shape in slide.shapes:
        text = ""
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = ''.join( run.text.strip() ) + ' '
                if text != '':
                    cleaned_text += text
    text_runs.append(cleaned_text.lstrip())
                     
for run in text_runs:
    print( run )

print( "Total Slides: " + str(i-1) )
