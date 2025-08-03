from pptx import Presentation                                                                                                                                                              
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
import os
import shutil
import sys
import time
import datetime

Grouped_Triggers = ['umh', 'fws', 'scripture reading', 'prayer for illumination', "the lord’s prayer", "call to worship", "affirmation of faith"] 
Standalone_Triggers = ['passing of the peace', 'rev.', 'prelude', 'postlude', 'the children’s moment', 'offering our gifts', 'offertory', 'sending forth', 'proclamation of god’s word', 'benediction'] 

def get_sunday_date():
    today = datetime.date.today()
    weekday = today.weekday()  # 0 = Monday, 6 = Sunday

    # Calculate days until next Sunday
    days_to_sunday = 6 - weekday

    # Get Sunday's date
    sunday_date = today + datetime.timedelta(days=days_to_sunday)

    return sunday_date

def create_slide( prs, layout, category, content ):
    ''' create_slide creates a slide in a given format, with a given layout, using a given conent string, and adds it to a given presentation object'''
    slide = prs.slides.add_slide(layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 255, 0)

    logo_img = '../assets/CrossFlame_WhiteRed.png'
    textbox_img = '../assets/text_plate.png'
    
    # add text plate
    X = Inches(1.5)
    Y = Inches(6.5)
    height = Inches(2.5)
    width = Inches(14.5)
    textbox = slide.shapes.add_picture(textbox_img, X, Y, width=width, height=height)

    # add logo
    X = Inches(-1.25)
    Y = Inches(6.5)
    height = Inches(2.5)
    logo = slide.shapes.add_picture(logo_img, X, Y, height=height)

    # add text
    X = Inches(1.5)
    Y = Inches(6.5)
    height = Inches(2.5)
    width = Inches(14.5)
    text_box = slide.shapes.add_textbox(X, Y, width, height)
    text_frame = text_box.text_frame
    text_frame.auto_size = None
    text_frame.word_wrap = True
    text_frame.text = content

    # Center the text horizontally in the text box
    text = text_frame.paragraphs[0]  # Access the first paragraph (text content)
    text.alignment = PP_ALIGN.CENTER  # Center the text horizontally

    # Center the text vertically in the text box
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center the text vertically

    font = text.font
    font.name = 'Arial'  # Set font name
    if category == "Standalone":
        font.size = Pt(40)
        font.bold = True
    elif category == "Grouped":
        font.size = Pt(32)
        font.bold = True

def extract_data( prs ):
    '''extract_data takes a given Presentation object and returns a list of lists containing the slide number, slide "type", and slide content of that object'''
    data = []
    Flag = "Standalone"
    slide_no = 1

    counter = 0
    in_sequence = False

    for slide in prs.slides:
        
        raw_data = []
        raw_data.append( slide_no )
        cleaned_text = ""

        for shape in slide.shapes:
            raw_text = ""

            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    raw_text = ''.join( run.text.strip() ) + ' '
                    if any( substr in raw_text.lower() for substr in Grouped_Triggers) == True:
                        Flag = "Grouped"
                        in_sequence == True
                        counter = 0
                    if any( substr in raw_text.lower() for substr in Standalone_Triggers) == True:
                        Flag = "Standalone"
                        counter = 0
                        in_sequence == False
                    if raw_text != '':
                        cleaned_text += raw_text

        if Flag == "Grouped":
            counter += 1
        else:
            in_sequence = False
            counter = 0
        #print( Flag )
        #print( in_sequence)
        #print( counter )

        raw_data.append( Flag )
        raw_data.append( cleaned_text.lstrip() )
        data.append( raw_data )
        slide_no += 1

        with open("ECOG_Input.txt", "a") as f:
            f.write(str(counter)+"\n")
        f.close()

    return data

def Main():
    if len( sys.argv ) != 2:
        if len( sys.argv ) < 2:
            print( "Error: No path to .pptx file provided" )
            exit( 1 )
        elif len( sys.argv ) > 2:
            print( "Error: To many arguments provided [Expected 1 argument, " + str( len( sys.argv ) - 1 ) + " arguments provided]" )
            exit( 1 )
    else:
        file_path = sys.argv[1]
        if file_path[-5::1] != ".pptx":
            print( "Error: expected .pptx file" )
            exit( 1 )
        else:
            template = Presentation( file_path )
            presentation = Presentation()
            presentation.slide_width = Inches(16)  # Width in inches
            presentation.slide_height = Inches(9)  # Height in inches
            blank_slide_layout = presentation.slide_layouts[6]
            
            date = get_sunday_date()
            os.makedirs(str(date))
            shutil.copy(file_path, str(date))
            os.chdir(str(date))

            macro_file = '../assets/ECOG_macro.txt'
            with open( macro_file, "r" ) as f:
                content = f.read()
            content = content.replace( "[insert date here]", str(date))
            output = f"ECOG_macro.txt"
            with open( output, "w" ) as f:
                f.write(content)

            data = extract_data( template )

            for raw_data in data:
                slide_no = raw_data[0]
                slide_type = raw_data[1]
                slide_content = raw_data[2]

                create_slide( presentation, blank_slide_layout, slide_type, slide_content )

        presentation.save( str(sys.argv[1])[:-5] + '-LT.pptx' )

if __name__ == "__main__":
    start_time = time.time()
    print("...Creating Lower Thirds...")
    Main()
    print("Lower Thirds creation complete!")
    print("Completed in --- %.5f seconds ---" % (time.time() - start_time))
